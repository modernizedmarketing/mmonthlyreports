#!/usr/bin/env python3
"""
Tag a finished monthly report PPTX with {{PLACEHOLDER}} tokens so it can be
reused as a template for the next month.

Usage:
    python3 tools/tag_template.py \
        --input  "data/Funded Profit/2026-03/previous_report.pptx" \
        --output "data/Funded Profit/2026-04/previous_report.pptx" \
        --client "One Funded" \
        --month  March \
        --year   2026

The script replaces the concrete values from the source report (e.g. "One Funded",
"March 2026", "$9,194.94") with their corresponding {{TOKEN}} strings while
preserving all font formatting (name, size, bold, colour).
"""
from __future__ import annotations

import argparse
import re
import shutil
import sys
from pathlib import Path
from typing import Dict, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor


# ─── formatting helpers (mirror fill_pptx.py) ────────────────────────────────

def _fmt_usd(v: float) -> str:
    return f"${v:,.2f}"

def _fmt_pct(v: float) -> str:
    return f"{v:.2f}%"

def _fmt_int(v: float) -> str:
    return f"{int(v):,}"

def _fmt_roas(v: float) -> str:
    return f"{v:.2f}"


# ─── low-level run replacement (mirrors fill_pptx.py) ────────────────────────

def _safe_rgb(run):
    try:
        return run.font.color.rgb
    except Exception:
        return None


def _replace_in_run(run, replacements: dict) -> None:
    text = run.text
    for old, new in replacements.items():
        if old in text:
            text = text.replace(old, new)
    if text != run.text:
        font_name = run.font.name
        font_size = run.font.size
        bold      = run.font.bold
        rgb       = _safe_rgb(run)
        run.text = text
        if font_name:        run.font.name  = font_name
        if font_size:        run.font.size  = font_size
        if bold is not None: run.font.bold  = bold
        if rgb:              run.font.color.rgb = rgb


def _replace_in_shape(shape, replacements: dict) -> None:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                _replace_in_run(run, replacements)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        _replace_in_run(run, replacements)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            _replace_in_shape(child, replacements)


# ─── build the old-value → token replacement map ─────────────────────────────

def build_tag_map(
    client: str,
    month: str,
    year: int,
    prev_month: str,
    next_month: str,
    # ── Slide 4 Total Ads ──────────────────────────────────────────────
    slide4_ad_revenue: float = 0,
    slide4_ad_cost: float = 0,
    slide4_roas: float = 0,
    slide4_cps: float = 0,
    slide4_l2s: float = 0,
    slide4_cvr: float = 0,
    slide4_aov: float = 0,
    company_revenue: float = 0,
    pct_rev_from_ads: float = 0,
    # ── Slide 4 previous month ─────────────────────────────────────────
    slide4_ad_revenue_prev: float = 0,
    slide4_ad_cost_prev: float = 0,
    slide4_roas_prev: float = 0,
    slide4_cps_prev: float = 0,
    slide4_l2s_prev: float = 0,
    slide4_cvr_prev: float = 0,
    slide4_aov_prev: float = 0,
    company_revenue_prev: float = 0,
    pct_rev_from_ads_prev: float = 0,
    # ── Slide 6 Google KPIs ────────────────────────────────────────────
    google_roas: float = 0,
    google_cps: float = 0,
    google_l2s: float = 0,
    google_cvr: float = 0,
    google_ctr: float = 0,
    google_revenue: float = 0,
    google_cost: float = 0,
    google_sales: float = 0,
    google_roas_prev: float = 0,
    google_cps_prev: float = 0,
    google_l2s_prev: float = 0,
    google_cvr_prev: float = 0,
    google_ctr_prev: float = 0,
    # ── Slide 8 Google volume ──────────────────────────────────────────
    google_clicks: float = 0,
    google_impressions: float = 0,
    google_leads: float = 0,
    # ── Slide 9 Google funnel ──────────────────────────────────────────
    google_tof_revenue: float = 0, google_tof_roas: float = 0, google_tof_sales: float = 0, google_tof_cps: float = 0, google_tof_cr: float = 0,
    google_mof_revenue: float = 0, google_mof_roas: float = 0, google_mof_sales: float = 0, google_mof_cps: float = 0, google_mof_cr: float = 0,
    google_bof_revenue: float = 0, google_bof_roas: float = 0, google_bof_sales: float = 0, google_bof_cps: float = 0, google_bof_cr: float = 0,
    # ── Slide 12 Meta KPIs ─────────────────────────────────────────────
    meta_roas: float = 0,
    meta_cps: float = 0,
    meta_l2s: float = 0,
    meta_cvr: float = 0,
    meta_ctr: float = 0,
    meta_revenue: float = 0,
    meta_cost: float = 0,
    meta_sales: float = 0,
    meta_roas_prev: float = 0,
    meta_cps_prev: float = 0,
    meta_l2s_prev: float = 0,
    meta_cvr_prev: float = 0,
    meta_ctr_prev: float = 0,
    # ── Slide 14 Meta volume ───────────────────────────────────────────
    meta_clicks: float = 0,
    meta_impressions: float = 0,
    meta_leads: float = 0,
    # ── Slides 15/16 Meta funnel ───────────────────────────────────────
    meta_tof_revenue: float = 0, meta_tof_roas: float = 0, meta_tof_sales: float = 0, meta_tof_cps: float = 0, meta_tof_cr: float = 0,
    meta_mof_revenue: float = 0, meta_mof_roas: float = 0, meta_mof_sales: float = 0, meta_mof_cps: float = 0, meta_mof_cr: float = 0,
    meta_bof_revenue: float = 0, meta_bof_roas: float = 0, meta_bof_sales: float = 0, meta_bof_cps: float = 0, meta_bof_cr: float = 0,
    # ── Slide 19 Bing KPIs ─────────────────────────────────────────────
    bing_roas: float = 0, bing_cps: float = 0, bing_l2s: float = 0, bing_cvr: float = 0, bing_ctr: float = 0,
    bing_revenue: float = 0, bing_cost: float = 0, bing_sales: float = 0,
    bing_roas_prev: float = 0, bing_cps_prev: float = 0, bing_l2s_prev: float = 0, bing_cvr_prev: float = 0, bing_ctr_prev: float = 0,
    # ── Slide 21 Bing volume ───────────────────────────────────────────
    bing_clicks: float = 0, bing_impressions: float = 0, bing_leads: float = 0,
    # ── Slide 22 Bing funnel ───────────────────────────────────────────
    bing_tof_revenue: float = 0, bing_tof_roas: float = 0, bing_tof_sales: float = 0, bing_tof_cps: float = 0, bing_tof_cr: float = 0,
    bing_mof_revenue: float = 0, bing_mof_roas: float = 0, bing_mof_sales: float = 0, bing_mof_cps: float = 0, bing_mof_cr: float = 0,
) -> Dict[str, str]:
    """
    Returns {old_value_string: "{{PLACEHOLDER}}"} for every dynamic field.
    Only entries where old_value is non-empty are included to avoid
    accidentally replacing "0" or "$0.00" across many slides.
    """
    raw: dict[str, str] = {
        # Slide 1
        client:                      "{{CLIENT}}",
        f"{month} {year}":           "{{MONTH}} {{YEAR}}",
        # Slide 4 — current
        _fmt_usd(slide4_ad_revenue): "{{SLIDE4_AD_REVENUE}}",
        _fmt_usd(slide4_ad_cost):    "{{SLIDE4_AD_COST}}",
        _fmt_roas(slide4_roas):      "{{SLIDE4_ROAS}}",
        _fmt_usd(slide4_cps):        "{{SLIDE4_CPS}}",
        _fmt_pct(slide4_l2s):        "{{SLIDE4_L2S}}",
        _fmt_pct(slide4_cvr):        "{{SLIDE4_CVR}}",
        _fmt_usd(slide4_aov):        "{{SLIDE4_AOV}}",
        _fmt_usd(company_revenue):   "{{COMPANY_REVENUE}}",
        _fmt_pct(pct_rev_from_ads):  "{{PCT_REV_FROM_ADS}}",
        # Slide 4 — previous
        _fmt_usd(slide4_ad_revenue_prev): "{{SLIDE4_AD_REVENUE_PREV}}",
        _fmt_usd(slide4_ad_cost_prev):    "{{SLIDE4_AD_COST_PREV}}",
        _fmt_roas(slide4_roas_prev):      "{{SLIDE4_ROAS_PREV}}",
        _fmt_usd(slide4_cps_prev):        "{{SLIDE4_CPS_PREV}}",
        _fmt_pct(slide4_l2s_prev):        "{{SLIDE4_L2S_PREV}}",
        _fmt_pct(slide4_cvr_prev):        "{{SLIDE4_CVR_PREV}}",
        _fmt_usd(slide4_aov_prev):        "{{SLIDE4_AOV_PREV}}",
        _fmt_usd(company_revenue_prev):   "{{COMPANY_REVENUE_PREV}}",
        _fmt_pct(pct_rev_from_ads_prev):  "{{PCT_REV_FROM_ADS_PREV}}",
        # Slide 6 — Google current
        _fmt_roas(google_roas):      "{{GOOGLE_ROAS}}",
        _fmt_usd(google_cps):        "{{GOOGLE_CPS}}",
        _fmt_pct(google_l2s):        "{{GOOGLE_L2S}}",
        _fmt_pct(google_cvr):        "{{GOOGLE_CVR}}",
        _fmt_pct(google_ctr):        "{{GOOGLE_CTR}}",
        _fmt_usd(google_revenue):    "{{GOOGLE_REVENUE}}",
        _fmt_usd(google_cost):       "{{GOOGLE_COST}}",
        _fmt_int(google_sales):      "{{GOOGLE_SALES}}",
        # Slide 6 — Google previous
        _fmt_roas(google_roas_prev): "{{GOOGLE_ROAS_PREV}}",
        _fmt_usd(google_cps_prev):   "{{GOOGLE_CPS_PREV}}",
        _fmt_pct(google_l2s_prev):   "{{GOOGLE_L2S_PREV}}",
        _fmt_pct(google_cvr_prev):   "{{GOOGLE_CVR_PREV}}",
        _fmt_pct(google_ctr_prev):   "{{GOOGLE_CTR_PREV}}",
        # Slide 8 — Google volume
        _fmt_int(google_clicks):      "{{GOOGLE_CLICKS}}",
        _fmt_int(google_impressions): "{{GOOGLE_IMPRESSIONS}}",
        _fmt_int(google_leads):       "{{GOOGLE_LEADS}}",
        # Slide 9 — Google funnel
        _fmt_usd(google_tof_revenue): "{{GOOGLE_TOF_REVENUE}}", _fmt_roas(google_tof_roas): "{{GOOGLE_TOF_ROAS}}", _fmt_int(google_tof_sales): "{{GOOGLE_TOF_SALES}}", _fmt_usd(google_tof_cps): "{{GOOGLE_TOF_CPS}}", _fmt_pct(google_tof_cr): "{{GOOGLE_TOF_CR}}",
        _fmt_usd(google_mof_revenue): "{{GOOGLE_MOF_REVENUE}}", _fmt_roas(google_mof_roas): "{{GOOGLE_MOF_ROAS}}", _fmt_int(google_mof_sales): "{{GOOGLE_MOF_SALES}}", _fmt_usd(google_mof_cps): "{{GOOGLE_MOF_CPS}}", _fmt_pct(google_mof_cr): "{{GOOGLE_MOF_CR}}",
        _fmt_usd(google_bof_revenue): "{{GOOGLE_BOF_REVENUE}}", _fmt_roas(google_bof_roas): "{{GOOGLE_BOF_ROAS}}", _fmt_int(google_bof_sales): "{{GOOGLE_BOF_SALES}}", _fmt_usd(google_bof_cps): "{{GOOGLE_BOF_CPS}}", _fmt_pct(google_bof_cr): "{{GOOGLE_BOF_CR}}",
        # Slide 12 — Meta current
        _fmt_roas(meta_roas):        "{{META_ROAS}}",
        _fmt_usd(meta_cps):          "{{META_CPS}}",
        _fmt_pct(meta_l2s):          "{{META_L2S}}",
        _fmt_pct(meta_cvr):          "{{META_CVR}}",
        _fmt_pct(meta_ctr):          "{{META_CTR}}",
        _fmt_usd(meta_revenue):      "{{META_REVENUE}}",
        _fmt_usd(meta_cost):         "{{META_COST}}",
        _fmt_int(meta_sales):        "{{META_SALES}}",
        # Slide 12 — Meta previous
        _fmt_roas(meta_roas_prev):   "{{META_ROAS_PREV}}",
        _fmt_usd(meta_cps_prev):     "{{META_CPS_PREV}}",
        _fmt_pct(meta_l2s_prev):     "{{META_L2S_PREV}}",
        _fmt_pct(meta_cvr_prev):     "{{META_CVR_PREV}}",
        _fmt_pct(meta_ctr_prev):     "{{META_CTR_PREV}}",
        # Slide 14 — Meta volume
        _fmt_int(meta_clicks):       "{{META_CLICKS}}",
        _fmt_int(meta_impressions):  "{{META_IMPRESSIONS}}",
        _fmt_int(meta_leads):        "{{META_LEADS}}",
        # Slides 15/16 — Meta funnel
        _fmt_usd(meta_tof_revenue): "{{META_TOF_REVENUE}}", _fmt_roas(meta_tof_roas): "{{META_TOF_ROAS}}", _fmt_int(meta_tof_sales): "{{META_TOF_SALES}}", _fmt_usd(meta_tof_cps): "{{META_TOF_CPS}}", _fmt_pct(meta_tof_cr): "{{META_TOF_CR}}",
        _fmt_usd(meta_mof_revenue): "{{META_MOF_REVENUE}}", _fmt_roas(meta_mof_roas): "{{META_MOF_ROAS}}", _fmt_int(meta_mof_sales): "{{META_MOF_SALES}}", _fmt_usd(meta_mof_cps): "{{META_MOF_CPS}}", _fmt_pct(meta_mof_cr): "{{META_MOF_CR}}",
        _fmt_usd(meta_bof_revenue): "{{META_BOF_REVENUE}}", _fmt_roas(meta_bof_roas): "{{META_BOF_ROAS}}", _fmt_int(meta_bof_sales): "{{META_BOF_SALES}}", _fmt_usd(meta_bof_cps): "{{META_BOF_CPS}}", _fmt_pct(meta_bof_cr): "{{META_BOF_CR}}",
        # Slide 19 — Bing current
        _fmt_roas(bing_roas): "{{BING_ROAS}}", _fmt_usd(bing_cps): "{{BING_CPS}}", _fmt_pct(bing_l2s): "{{BING_L2S}}", _fmt_pct(bing_cvr): "{{BING_CVR}}", _fmt_pct(bing_ctr): "{{BING_CTR}}",
        _fmt_usd(bing_revenue): "{{BING_REVENUE}}", _fmt_usd(bing_cost): "{{BING_COST}}", _fmt_int(bing_sales): "{{BING_SALES}}",
        # Slide 19 — Bing previous
        _fmt_roas(bing_roas_prev): "{{BING_ROAS_PREV}}", _fmt_usd(bing_cps_prev): "{{BING_CPS_PREV}}", _fmt_pct(bing_l2s_prev): "{{BING_L2S_PREV}}", _fmt_pct(bing_cvr_prev): "{{BING_CVR_PREV}}", _fmt_pct(bing_ctr_prev): "{{BING_CTR_PREV}}",
        # Slide 21 — Bing volume
        _fmt_int(bing_clicks): "{{BING_CLICKS}}", _fmt_int(bing_impressions): "{{BING_IMPRESSIONS}}", _fmt_int(bing_leads): "{{BING_LEADS}}",
        # Slide 22 — Bing funnel
        _fmt_usd(bing_tof_revenue): "{{BING_TOF_REVENUE}}", _fmt_roas(bing_tof_roas): "{{BING_TOF_ROAS}}", _fmt_int(bing_tof_sales): "{{BING_TOF_SALES}}", _fmt_usd(bing_tof_cps): "{{BING_TOF_CPS}}", _fmt_pct(bing_tof_cr): "{{BING_TOF_CR}}",
        _fmt_usd(bing_mof_revenue): "{{BING_MOF_REVENUE}}", _fmt_roas(bing_mof_roas): "{{BING_MOF_ROAS}}", _fmt_int(bing_mof_sales): "{{BING_MOF_SALES}}", _fmt_usd(bing_mof_cps): "{{BING_MOF_CPS}}", _fmt_pct(bing_mof_cr): "{{BING_MOF_CR}}",
    }

    # Drop keys that are zero/empty to avoid replacing generic "0" or "$0.00"
    _skip = {"$0.00", "0", "0.00", "0.00%", "0,"}
    return {k: v for k, v in raw.items() if k and k not in _skip}


def tag_pptx(input_path: Path, output_path: Path, tag_map: dict) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(input_path, output_path)
    prs = Presentation(output_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_in_shape(shape, tag_map)
    prs.save(output_path)


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Tag a finished PPTX with {{PLACEHOLDER}} tokens.")
    p.add_argument("--input",  required=True, help="Path to last month's finished PPTX")
    p.add_argument("--output", required=True, help="Where to write the tagged template")
    p.add_argument("--client", required=True)
    p.add_argument("--month",  required=True, help="Month of the source report (e.g. March)")
    p.add_argument("--year",   required=True, type=int)
    p.add_argument("--prev-month", default="")
    p.add_argument("--next-month", default="")
    # Slide 4 total ads
    p.add_argument("--slide4-ad-revenue", type=float, default=0)
    p.add_argument("--slide4-ad-cost",    type=float, default=0)
    p.add_argument("--slide4-roas",       type=float, default=0)
    p.add_argument("--slide4-cps",        type=float, default=0)
    p.add_argument("--slide4-l2s",        type=float, default=0)
    p.add_argument("--slide4-cvr",        type=float, default=0)
    p.add_argument("--slide4-aov",        type=float, default=0)
    p.add_argument("--company-revenue",   type=float, default=0)
    p.add_argument("--pct-rev-from-ads",  type=float, default=0)
    # Slide 4 prev
    p.add_argument("--slide4-ad-revenue-prev", type=float, default=0)
    p.add_argument("--slide4-ad-cost-prev",    type=float, default=0)
    p.add_argument("--slide4-roas-prev",       type=float, default=0)
    p.add_argument("--slide4-cps-prev",        type=float, default=0)
    p.add_argument("--slide4-l2s-prev",        type=float, default=0)
    p.add_argument("--slide4-cvr-prev",        type=float, default=0)
    p.add_argument("--slide4-aov-prev",        type=float, default=0)
    p.add_argument("--company-revenue-prev",   type=float, default=0)
    p.add_argument("--pct-rev-from-ads-prev",  type=float, default=0)
    # Google
    p.add_argument("--google-roas",        type=float, default=0)
    p.add_argument("--google-cps",         type=float, default=0)
    p.add_argument("--google-l2s",         type=float, default=0)
    p.add_argument("--google-cvr",         type=float, default=0)
    p.add_argument("--google-ctr",         type=float, default=0)
    p.add_argument("--google-revenue",     type=float, default=0)
    p.add_argument("--google-cost",        type=float, default=0)
    p.add_argument("--google-sales",       type=float, default=0)
    p.add_argument("--google-roas-prev",   type=float, default=0)
    p.add_argument("--google-cps-prev",    type=float, default=0)
    p.add_argument("--google-l2s-prev",    type=float, default=0)
    p.add_argument("--google-cvr-prev",    type=float, default=0)
    p.add_argument("--google-ctr-prev",    type=float, default=0)
    p.add_argument("--google-clicks",      type=float, default=0)
    p.add_argument("--google-impressions", type=float, default=0)
    p.add_argument("--google-leads",       type=float, default=0)
    # Google funnel
    p.add_argument("--google-tof-revenue", type=float, default=0); p.add_argument("--google-tof-roas", type=float, default=0); p.add_argument("--google-tof-sales", type=float, default=0); p.add_argument("--google-tof-cps", type=float, default=0); p.add_argument("--google-tof-cr", type=float, default=0)
    p.add_argument("--google-mof-revenue", type=float, default=0); p.add_argument("--google-mof-roas", type=float, default=0); p.add_argument("--google-mof-sales", type=float, default=0); p.add_argument("--google-mof-cps", type=float, default=0); p.add_argument("--google-mof-cr", type=float, default=0)
    p.add_argument("--google-bof-revenue", type=float, default=0); p.add_argument("--google-bof-roas", type=float, default=0); p.add_argument("--google-bof-sales", type=float, default=0); p.add_argument("--google-bof-cps", type=float, default=0); p.add_argument("--google-bof-cr", type=float, default=0)
    # Meta
    p.add_argument("--meta-roas",        type=float, default=0)
    p.add_argument("--meta-cps",         type=float, default=0)
    p.add_argument("--meta-l2s",         type=float, default=0)
    p.add_argument("--meta-cvr",         type=float, default=0)
    p.add_argument("--meta-ctr",         type=float, default=0)
    p.add_argument("--meta-revenue",     type=float, default=0)
    p.add_argument("--meta-cost",        type=float, default=0)
    p.add_argument("--meta-sales",       type=float, default=0)
    p.add_argument("--meta-roas-prev",   type=float, default=0)
    p.add_argument("--meta-cps-prev",    type=float, default=0)
    p.add_argument("--meta-l2s-prev",    type=float, default=0)
    p.add_argument("--meta-cvr-prev",    type=float, default=0)
    p.add_argument("--meta-ctr-prev",    type=float, default=0)
    p.add_argument("--meta-clicks",      type=float, default=0)
    p.add_argument("--meta-impressions", type=float, default=0)
    p.add_argument("--meta-leads",       type=float, default=0)
    # Meta funnel
    p.add_argument("--meta-tof-revenue", type=float, default=0); p.add_argument("--meta-tof-roas", type=float, default=0); p.add_argument("--meta-tof-sales", type=float, default=0); p.add_argument("--meta-tof-cps", type=float, default=0); p.add_argument("--meta-tof-cr", type=float, default=0)
    p.add_argument("--meta-mof-revenue", type=float, default=0); p.add_argument("--meta-mof-roas", type=float, default=0); p.add_argument("--meta-mof-sales", type=float, default=0); p.add_argument("--meta-mof-cps", type=float, default=0); p.add_argument("--meta-mof-cr", type=float, default=0)
    p.add_argument("--meta-bof-revenue", type=float, default=0); p.add_argument("--meta-bof-roas", type=float, default=0); p.add_argument("--meta-bof-sales", type=float, default=0); p.add_argument("--meta-bof-cps", type=float, default=0); p.add_argument("--meta-bof-cr", type=float, default=0)
    # Bing
    p.add_argument("--bing-roas", type=float, default=0); p.add_argument("--bing-cps", type=float, default=0); p.add_argument("--bing-l2s", type=float, default=0); p.add_argument("--bing-cvr", type=float, default=0); p.add_argument("--bing-ctr", type=float, default=0)
    p.add_argument("--bing-revenue", type=float, default=0); p.add_argument("--bing-cost", type=float, default=0); p.add_argument("--bing-sales", type=float, default=0)
    p.add_argument("--bing-roas-prev", type=float, default=0); p.add_argument("--bing-cps-prev", type=float, default=0); p.add_argument("--bing-l2s-prev", type=float, default=0); p.add_argument("--bing-cvr-prev", type=float, default=0); p.add_argument("--bing-ctr-prev", type=float, default=0)
    p.add_argument("--bing-clicks", type=float, default=0); p.add_argument("--bing-impressions", type=float, default=0); p.add_argument("--bing-leads", type=float, default=0)
    p.add_argument("--bing-tof-revenue", type=float, default=0); p.add_argument("--bing-tof-roas", type=float, default=0); p.add_argument("--bing-tof-sales", type=float, default=0); p.add_argument("--bing-tof-cps", type=float, default=0); p.add_argument("--bing-tof-cr", type=float, default=0)
    p.add_argument("--bing-mof-revenue", type=float, default=0); p.add_argument("--bing-mof-roas", type=float, default=0); p.add_argument("--bing-mof-sales", type=float, default=0); p.add_argument("--bing-mof-cps", type=float, default=0); p.add_argument("--bing-mof-cr", type=float, default=0)
    return p.parse_args()


def main() -> int:
    args = parse_args()
    tag_map = build_tag_map(
        client=args.client, month=args.month, year=args.year,
        prev_month=args.prev_month, next_month=args.next_month,
        slide4_ad_revenue=args.slide4_ad_revenue, slide4_ad_cost=args.slide4_ad_cost,
        slide4_roas=args.slide4_roas, slide4_cps=args.slide4_cps,
        slide4_l2s=args.slide4_l2s, slide4_cvr=args.slide4_cvr, slide4_aov=args.slide4_aov,
        company_revenue=args.company_revenue, pct_rev_from_ads=args.pct_rev_from_ads,
        slide4_ad_revenue_prev=args.slide4_ad_revenue_prev, slide4_ad_cost_prev=args.slide4_ad_cost_prev,
        slide4_roas_prev=args.slide4_roas_prev, slide4_cps_prev=args.slide4_cps_prev,
        slide4_l2s_prev=args.slide4_l2s_prev, slide4_cvr_prev=args.slide4_cvr_prev, slide4_aov_prev=args.slide4_aov_prev,
        company_revenue_prev=args.company_revenue_prev, pct_rev_from_ads_prev=args.pct_rev_from_ads_prev,
        google_roas=args.google_roas, google_cps=args.google_cps, google_l2s=args.google_l2s,
        google_cvr=args.google_cvr, google_ctr=args.google_ctr, google_revenue=args.google_revenue,
        google_cost=args.google_cost, google_sales=args.google_sales,
        google_roas_prev=args.google_roas_prev, google_cps_prev=args.google_cps_prev,
        google_l2s_prev=args.google_l2s_prev, google_cvr_prev=args.google_cvr_prev, google_ctr_prev=args.google_ctr_prev,
        google_clicks=args.google_clicks, google_impressions=args.google_impressions, google_leads=args.google_leads,
        google_tof_revenue=args.google_tof_revenue, google_tof_roas=args.google_tof_roas, google_tof_sales=args.google_tof_sales, google_tof_cps=args.google_tof_cps, google_tof_cr=args.google_tof_cr,
        google_mof_revenue=args.google_mof_revenue, google_mof_roas=args.google_mof_roas, google_mof_sales=args.google_mof_sales, google_mof_cps=args.google_mof_cps, google_mof_cr=args.google_mof_cr,
        google_bof_revenue=args.google_bof_revenue, google_bof_roas=args.google_bof_roas, google_bof_sales=args.google_bof_sales, google_bof_cps=args.google_bof_cps, google_bof_cr=args.google_bof_cr,
        meta_roas=args.meta_roas, meta_cps=args.meta_cps, meta_l2s=args.meta_l2s,
        meta_cvr=args.meta_cvr, meta_ctr=args.meta_ctr, meta_revenue=args.meta_revenue,
        meta_cost=args.meta_cost, meta_sales=args.meta_sales,
        meta_roas_prev=args.meta_roas_prev, meta_cps_prev=args.meta_cps_prev,
        meta_l2s_prev=args.meta_l2s_prev, meta_cvr_prev=args.meta_cvr_prev, meta_ctr_prev=args.meta_ctr_prev,
        meta_clicks=args.meta_clicks, meta_impressions=args.meta_impressions, meta_leads=args.meta_leads,
        meta_tof_revenue=args.meta_tof_revenue, meta_tof_roas=args.meta_tof_roas, meta_tof_sales=args.meta_tof_sales, meta_tof_cps=args.meta_tof_cps, meta_tof_cr=args.meta_tof_cr,
        meta_mof_revenue=args.meta_mof_revenue, meta_mof_roas=args.meta_mof_roas, meta_mof_sales=args.meta_mof_sales, meta_mof_cps=args.meta_mof_cps, meta_mof_cr=args.meta_mof_cr,
        meta_bof_revenue=args.meta_bof_revenue, meta_bof_roas=args.meta_bof_roas, meta_bof_sales=args.meta_bof_sales, meta_bof_cps=args.meta_bof_cps, meta_bof_cr=args.meta_bof_cr,
        bing_roas=args.bing_roas, bing_cps=args.bing_cps, bing_l2s=args.bing_l2s, bing_cvr=args.bing_cvr, bing_ctr=args.bing_ctr,
        bing_revenue=args.bing_revenue, bing_cost=args.bing_cost, bing_sales=args.bing_sales,
        bing_roas_prev=args.bing_roas_prev, bing_cps_prev=args.bing_cps_prev, bing_l2s_prev=args.bing_l2s_prev, bing_cvr_prev=args.bing_cvr_prev, bing_ctr_prev=args.bing_ctr_prev,
        bing_clicks=args.bing_clicks, bing_impressions=args.bing_impressions, bing_leads=args.bing_leads,
        bing_tof_revenue=args.bing_tof_revenue, bing_tof_roas=args.bing_tof_roas, bing_tof_sales=args.bing_tof_sales, bing_tof_cps=args.bing_tof_cps, bing_tof_cr=args.bing_tof_cr,
        bing_mof_revenue=args.bing_mof_revenue, bing_mof_roas=args.bing_mof_roas, bing_mof_sales=args.bing_mof_sales, bing_mof_cps=args.bing_mof_cps, bing_mof_cr=args.bing_mof_cr,
    )

    print(f"Tagging {len(tag_map)} values → placeholders")
    tag_pptx(Path(args.input), Path(args.output), tag_map)
    print(f"Tagged template saved to: {args.output}")

    # Verify
    prs = Presentation(args.output)
    found: set = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                found.update(re.findall(r"\{\{[^}]+\}\}", shape.text_frame.text))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        found.update(re.findall(r"\{\{[^}]+\}\}", cell.text_frame.text))
    print(f"Placeholders written to template ({len(found)}):")
    for tok in sorted(found):
        print(f"  {tok}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
