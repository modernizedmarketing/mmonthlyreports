"""
Edits a PPTX template by replacing {{placeholder}} tokens with real values.
All font properties (name, size, bold, color) are preserved on replacement.
"""
import shutil
from pathlib import Path
from typing import Optional, Dict, Union
from pptx import Presentation
from pptx.dml.color import RGBColor


def _safe_rgb(run) -> Optional[RGBColor]:
    """Read font color without crashing on SchemeColor."""
    try:
        return run.font.color.rgb
    except Exception:
        return None


def _replace_in_run(run, replacements: dict):
    """Replace placeholder tokens in a single run, preserving font properties."""
    text = run.text
    for placeholder, value in replacements.items():
        if placeholder in text:
            text = text.replace(placeholder, str(value))
    if text != run.text:
        # Capture existing formatting
        font_name = run.font.name
        font_size = run.font.size
        bold      = run.font.bold
        rgb       = _safe_rgb(run)
        run.text = text
        # Restore formatting
        if font_name:       run.font.name = font_name
        if font_size:       run.font.size = font_size
        if bold is not None: run.font.bold = bold
        if rgb:             run.font.color.rgb = rgb


def _replace_in_shape(shape, replacements: dict):
    """Recursively replace placeholders in all text frames and tables."""
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                _replace_in_run(run, replacements)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        _replace_in_run(run, replacements)
    # Handle grouped shapes
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            _replace_in_shape(child, replacements)


def build_replacements(
    client: str,
    month: str,
    year: int,
    prev_month: str,
    next_month: str,
    kpis: dict,
    insights: dict,
    user_overrides: dict,
    prev_kpis: dict = None,
) -> dict:
    """Build the full placeholder → value dict for all slides.

    prev_kpis: optional KPI dict from last month's xlsx (enables _PREV placeholders).
    """
    g  = kpis["google"]
    m  = kpis["meta"]
    t  = kpis["totals"]
    gf = kpis["google_funnels"]
    mf = kpis["meta_funnels"]
    b  = kpis.get("bing", {})
    bf = kpis.get("bing_funnels", {})
    ov = user_overrides

    pg  = prev_kpis["google"]  if prev_kpis else {}
    pm  = prev_kpis["meta"]    if prev_kpis else {}
    pt  = prev_kpis["totals"]  if prev_kpis else {}
    pb  = prev_kpis.get("bing", {}) if prev_kpis else {}

    def fmt_pct(v): return f"{v:.2f}%" if v else "0.00%"
    def fmt_usd(v): return f"${v:,.2f}" if v else "$0.00"
    def fmt_int(v): return f"{int(v):,}" if v else "0"
    def fmt_roas(v): return f"{v:.2f}" if v else "0.00"
    def _p(d, key, fmt=str):
        """Safe prev-month value; empty string when prev data not loaded."""
        val = d.get(key)
        if val is None:
            return ""
        try:
            return fmt(val)
        except Exception:
            return str(val)

    def _funnel(fd, stage, key, fmt):
        return fmt(fd.get(stage, {}).get(key, 0)) if fd.get(stage) else ""

    ad_revenue  = ov.get("ad_revenue")  or t["revenue"]
    ad_cost     = ov.get("ad_cost")     or t["cost"]
    co_revenue  = ov.get("company_revenue", 0)
    slide4_roas = fmt_roas(ad_revenue / ad_cost) if ad_cost else fmt_roas(t["roas"])
    slide4_cps  = fmt_usd(ad_cost / t["sales"]) if t.get("sales") else "$0.00"
    slide4_l2s  = fmt_pct(t.get("l2s_pct", 0))
    slide4_cvr  = fmt_pct(t.get("cvr_pct", 0))
    slide4_aov  = fmt_usd(t.get("aov", 0))
    pct_rev     = fmt_pct(ad_revenue / co_revenue * 100) if co_revenue else "0.00%"

    p_ad_rev    = fmt_usd(ov.get("prev_ad_revenue") or pt.get("revenue", 0)) if prev_kpis else ""
    p_ad_cost   = fmt_usd(ov.get("prev_ad_cost")    or pt.get("cost", 0))    if prev_kpis else ""
    p_co_rev    = fmt_usd(ov.get("prev_company_revenue", 0)) if prev_kpis else ""
    p_roas      = fmt_roas(pt.get("roas", 0)) if prev_kpis else ""
    p_cps       = fmt_usd(pt.get("cps", 0))  if prev_kpis else ""
    p_l2s       = fmt_pct(pt.get("l2s_pct", 0)) if prev_kpis else ""
    p_cvr       = fmt_pct(pt.get("cvr_pct", 0)) if prev_kpis else ""
    p_aov       = fmt_usd(pt.get("aov", 0))  if prev_kpis else ""
    p_co_rev_val= ov.get("prev_company_revenue", 0) if prev_kpis else 0
    p_ad_rev_val= ov.get("prev_ad_revenue") or pt.get("revenue", 0) if prev_kpis else 0
    p_pct_rev   = fmt_pct(p_ad_rev_val / p_co_rev_val * 100) if prev_kpis and p_co_rev_val else ""

    action_items = insights.get("action_items", [])

    replacements = {
        # ── Slide 1 — Cover ──────────────────────────────────────────────────
        "{{MONTH}}":      month,
        "{{YEAR}}":       str(year),
        "{{CLIENT}}":     client,
        "{{PREV_MONTH}}": prev_month,
        "{{NEXT_MONTH}}": next_month,

        # ── Slide 4 — Total Ads (current month) ──────────────────────────────
        "{{SLIDE4_AD_REVENUE}}": fmt_usd(ad_revenue),
        "{{SLIDE4_AD_COST}}":    fmt_usd(ad_cost),
        "{{SLIDE4_ROAS}}":       slide4_roas,
        "{{SLIDE4_CPS}}":        slide4_cps,
        "{{SLIDE4_L2S}}":        slide4_l2s,
        "{{SLIDE4_CVR}}":        slide4_cvr,
        "{{SLIDE4_AOV}}":        slide4_aov,
        "{{COMPANY_REVENUE}}":   fmt_usd(co_revenue),
        "{{PCT_REV_FROM_ADS}}":  pct_rev,

        # ── Slide 4 — Total Ads (previous month) ─────────────────────────────
        "{{SLIDE4_AD_REVENUE_PREV}}":    p_ad_rev,
        "{{SLIDE4_AD_COST_PREV}}":       p_ad_cost,
        "{{SLIDE4_ROAS_PREV}}":          p_roas,
        "{{SLIDE4_CPS_PREV}}":           p_cps,
        "{{SLIDE4_L2S_PREV}}":           p_l2s,
        "{{SLIDE4_CVR_PREV}}":           p_cvr,
        "{{SLIDE4_AOV_PREV}}":           p_aov,
        "{{COMPANY_REVENUE_PREV}}":      p_co_rev,
        "{{PCT_REV_FROM_ADS_PREV}}":     p_pct_rev,

        # ── Slide 6 — Google KPIs (current) ──────────────────────────────────
        "{{GOOGLE_ROAS}}":        fmt_roas(g["roas"]),
        "{{GOOGLE_CPS}}":         fmt_usd(g["cps"]),
        "{{GOOGLE_L2S}}":         fmt_pct(g["l2s_pct"]),
        "{{GOOGLE_CVR}}":         fmt_pct(g["cvr_pct"]),
        "{{GOOGLE_CTR}}":         fmt_pct(g["ctr_pct"]),
        "{{GOOGLE_REVENUE}}":     fmt_usd(g["revenue"]),
        "{{GOOGLE_COST}}":        fmt_usd(g["cost"]),
        "{{GOOGLE_SALES}}":       fmt_int(g["sales"]),

        # ── Slide 6 — Google KPIs (previous) ─────────────────────────────────
        "{{GOOGLE_ROAS_PREV}}": _p(pg, "roas", fmt_roas),
        "{{GOOGLE_CPS_PREV}}":  _p(pg, "cps",  fmt_usd),
        "{{GOOGLE_L2S_PREV}}":  _p(pg, "l2s_pct", fmt_pct),
        "{{GOOGLE_CVR_PREV}}":  _p(pg, "cvr_pct", fmt_pct),
        "{{GOOGLE_CTR_PREV}}":  _p(pg, "ctr_pct", fmt_pct),

        # ── Slide 8 — Google volume ───────────────────────────────────────────
        "{{GOOGLE_CLICKS}}":      fmt_int(g["clicks"]),
        "{{GOOGLE_IMPRESSIONS}}": fmt_int(g["impressions"]),
        "{{GOOGLE_LEADS}}":       fmt_int(g["leads"]),

        # ── Slide 9 — Google funnel ───────────────────────────────────────────
        "{{GOOGLE_TOF_REVENUE}}": fmt_usd(gf.get("TOF", {}).get("revenue", 0)),
        "{{GOOGLE_TOF_COST}}":    fmt_usd(gf.get("TOF", {}).get("cost", 0)),
        "{{GOOGLE_TOF_ROAS}}":    fmt_roas(gf.get("TOF", {}).get("roas", 0)),
        "{{GOOGLE_TOF_SALES}}":   fmt_int(gf.get("TOF", {}).get("sales", 0)),
        "{{GOOGLE_TOF_CPS}}":     _funnel(gf, "TOF", "cps", fmt_usd),
        "{{GOOGLE_TOF_CR}}":      _funnel(gf, "TOF", "cvr_pct", fmt_pct),
        "{{GOOGLE_MOF_REVENUE}}": fmt_usd(gf.get("MOF", {}).get("revenue", 0)),
        "{{GOOGLE_MOF_COST}}":    fmt_usd(gf.get("MOF", {}).get("cost", 0)),
        "{{GOOGLE_MOF_ROAS}}":    fmt_roas(gf.get("MOF", {}).get("roas", 0)),
        "{{GOOGLE_MOF_SALES}}":   fmt_int(gf.get("MOF", {}).get("sales", 0)),
        "{{GOOGLE_MOF_CPS}}":     _funnel(gf, "MOF", "cps", fmt_usd),
        "{{GOOGLE_MOF_CR}}":      _funnel(gf, "MOF", "cvr_pct", fmt_pct),
        "{{GOOGLE_BOF_REVENUE}}": fmt_usd(gf.get("BOF", {}).get("revenue", 0)),
        "{{GOOGLE_BOF_COST}}":    fmt_usd(gf.get("BOF", {}).get("cost", 0)),
        "{{GOOGLE_BOF_ROAS}}":    fmt_roas(gf.get("BOF", {}).get("roas", 0)),
        "{{GOOGLE_BOF_SALES}}":   fmt_int(gf.get("BOF", {}).get("sales", 0)),
        "{{GOOGLE_BOF_CPS}}":     _funnel(gf, "BOF", "cps", fmt_usd),
        "{{GOOGLE_BOF_CR}}":      _funnel(gf, "BOF", "cvr_pct", fmt_pct),

        # ── Slides 12 — Meta KPIs (current) ──────────────────────────────────
        "{{META_ROAS}}":    fmt_roas(m["roas"]),
        "{{META_CPS}}":     fmt_usd(m["cps"]),
        "{{META_L2S}}":     fmt_pct(m["l2s_pct"]),
        "{{META_CVR}}":     fmt_pct(m["cvr_pct"]),
        "{{META_CTR}}":     fmt_pct(m["ctr_pct"]),
        "{{META_REVENUE}}": fmt_usd(m["revenue"]),
        "{{META_COST}}":    fmt_usd(m["cost"]),
        "{{META_SALES}}":   fmt_int(m["sales"]),

        # ── Slide 12 — Meta KPIs (previous) ──────────────────────────────────
        "{{META_ROAS_PREV}}": _p(pm, "roas", fmt_roas),
        "{{META_CPS_PREV}}":  _p(pm, "cps",  fmt_usd),
        "{{META_L2S_PREV}}":  _p(pm, "l2s_pct", fmt_pct),
        "{{META_CVR_PREV}}":  _p(pm, "cvr_pct", fmt_pct),
        "{{META_CTR_PREV}}":  _p(pm, "ctr_pct", fmt_pct),

        # ── Slide 14 — Meta volume ────────────────────────────────────────────
        "{{META_CLICKS}}":      fmt_int(m["clicks"]),
        "{{META_IMPRESSIONS}}": fmt_int(m["impressions"]),
        "{{META_LEADS}}":       fmt_int(m["leads"]),

        # ── Slides 15/16 — Meta funnel ────────────────────────────────────────
        "{{META_TOF_REVENUE}}": fmt_usd(mf.get("TOF", {}).get("revenue", 0)),
        "{{META_TOF_ROAS}}":    fmt_roas(mf.get("TOF", {}).get("roas", 0)),
        "{{META_TOF_SALES}}":   fmt_int(mf.get("TOF", {}).get("sales", 0)),
        "{{META_TOF_CPS}}":     _funnel(mf, "TOF", "cps", fmt_usd),
        "{{META_TOF_CR}}":      _funnel(mf, "TOF", "cvr_pct", fmt_pct),
        "{{META_MOF_REVENUE}}": fmt_usd(mf.get("MOF", {}).get("revenue", 0)),
        "{{META_MOF_ROAS}}":    fmt_roas(mf.get("MOF", {}).get("roas", 0)),
        "{{META_MOF_SALES}}":   fmt_int(mf.get("MOF", {}).get("sales", 0)),
        "{{META_MOF_CPS}}":     _funnel(mf, "MOF", "cps", fmt_usd),
        "{{META_MOF_CR}}":      _funnel(mf, "MOF", "cvr_pct", fmt_pct),
        "{{META_BOF_REVENUE}}": fmt_usd(mf.get("BOF", {}).get("revenue", 0)),
        "{{META_BOF_ROAS}}":    fmt_roas(mf.get("BOF", {}).get("roas", 0)),
        "{{META_BOF_SALES}}":   fmt_int(mf.get("BOF", {}).get("sales", 0)),
        "{{META_BOF_CPS}}":     _funnel(mf, "BOF", "cps", fmt_usd),
        "{{META_BOF_CR}}":      _funnel(mf, "BOF", "cvr_pct", fmt_pct),

        # ── Slide 19 — Bing KPIs (current) ───────────────────────────────────
        "{{BING_ROAS}}": fmt_roas(b.get("roas", 0)) if b else "",
        "{{BING_CPS}}":  fmt_usd(b.get("cps", 0))  if b else "",
        "{{BING_L2S}}":  fmt_pct(b.get("l2s_pct", 0)) if b else "",
        "{{BING_CVR}}":  fmt_pct(b.get("cvr_pct", 0)) if b else "",
        "{{BING_CTR}}":  fmt_pct(b.get("ctr_pct", 0)) if b else "",
        "{{BING_REVENUE}}": fmt_usd(b.get("revenue", 0)) if b else "",
        "{{BING_COST}}":    fmt_usd(b.get("cost", 0))    if b else "",
        "{{BING_SALES}}":   fmt_int(b.get("sales", 0))   if b else "",

        # ── Slide 19 — Bing KPIs (previous) ──────────────────────────────────
        "{{BING_ROAS_PREV}}": _p(pb, "roas", fmt_roas),
        "{{BING_CPS_PREV}}":  _p(pb, "cps",  fmt_usd),
        "{{BING_L2S_PREV}}":  _p(pb, "l2s_pct", fmt_pct),
        "{{BING_CVR_PREV}}":  _p(pb, "cvr_pct", fmt_pct),
        "{{BING_CTR_PREV}}":  _p(pb, "ctr_pct", fmt_pct),

        # ── Slide 21 — Bing volume ────────────────────────────────────────────
        "{{BING_CLICKS}}":      fmt_int(b.get("clicks", 0))      if b else "",
        "{{BING_IMPRESSIONS}}": fmt_int(b.get("impressions", 0)) if b else "",
        "{{BING_LEADS}}":       fmt_int(b.get("leads", 0))       if b else "",

        # ── Slide 22 — Bing funnel ────────────────────────────────────────────
        "{{BING_TOF_REVENUE}}": fmt_usd(bf.get("TOF", {}).get("revenue", 0)) if bf else "",
        "{{BING_TOF_ROAS}}":    fmt_roas(bf.get("TOF", {}).get("roas", 0))   if bf else "",
        "{{BING_TOF_SALES}}":   fmt_int(bf.get("TOF", {}).get("sales", 0))   if bf else "",
        "{{BING_TOF_CPS}}":     _funnel(bf, "TOF", "cps", fmt_usd)           if bf else "",
        "{{BING_TOF_CR}}":      _funnel(bf, "TOF", "cvr_pct", fmt_pct)       if bf else "",
        "{{BING_MOF_REVENUE}}": fmt_usd(bf.get("MOF", {}).get("revenue", 0)) if bf else "",
        "{{BING_MOF_ROAS}}":    fmt_roas(bf.get("MOF", {}).get("roas", 0))   if bf else "",
        "{{BING_MOF_SALES}}":   fmt_int(bf.get("MOF", {}).get("sales", 0))   if bf else "",
        "{{BING_MOF_CPS}}":     _funnel(bf, "MOF", "cps", fmt_usd)           if bf else "",
        "{{BING_MOF_CR}}":      _funnel(bf, "MOF", "cvr_pct", fmt_pct)       if bf else "",

        # ── Slide 3 — AI Insights ─────────────────────────────────────────────
        "{{SLIDE3_GENERAL_INSIGHTS}}": insights.get("slide3_general_insights", ""),
        "{{SLIDE3_BUDGET_ROAS}}":      insights.get("slide3_budget_roas", ""),
        "{{SLIDE3_STRATEGY}}":         insights.get("slide3_strategy", ""),
        "{{GOOGLE_TOP_PERFORMER}}":    insights.get("google_top_performer", ""),
        "{{GOOGLE_MAIN_DROP}}":        insights.get("google_main_drop", ""),
        "{{GOOGLE_NEXT_STEPS}}":       insights.get("google_next_steps", ""),
        "{{META_TOP_PERFORMER}}":      insights.get("meta_top_performer", ""),
        "{{META_MAIN_DROP}}":          insights.get("meta_main_drop", ""),
        "{{META_NEXT_STEPS}}":         insights.get("meta_next_steps", ""),
        "{{PM_NARRATIVE}}":            insights.get("performance_manager_narrative", ""),
        "{{ACTION_ITEM_1}}": action_items[0] if len(action_items) > 0 else "",
        "{{ACTION_ITEM_2}}": action_items[1] if len(action_items) > 1 else "",
        "{{ACTION_ITEM_3}}": action_items[2] if len(action_items) > 2 else "",
        "{{ACTION_ITEM_4}}": action_items[3] if len(action_items) > 3 else "",
        "{{ACTION_ITEM_5}}": action_items[4] if len(action_items) > 4 else "",
    }

    return replacements


def fill_pptx(
    template_path: Union[str, Path],
    output_path: Union[str, Path],
    replacements: Dict,
) -> Path:
    """
    Copy template PPTX to output_path, then replace all {{placeholder}} tokens.
    Returns the output path.
    """
    template_path = Path(template_path)
    output_path   = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    shutil.copy2(template_path, output_path)
    prs = Presentation(output_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_in_shape(shape, replacements)

    prs.save(output_path)
    return output_path
