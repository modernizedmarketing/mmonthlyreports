#!/usr/bin/env python3
"""Run the monthly marketing report pipeline against real local files."""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Iterable

from pptx import Presentation

if __package__ is None or __package__ == "":
    sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from tools.calculate_kpis import build_full_kpi_report
from tools.fill_pptx import build_replacements, fill_pptx
from tools.generate_insights import generate_insights
from tools.validate_data import load_sheet, validate_or_raise


REQUIRED_INSIGHT_KEYS = {
    "slide3_general_insights",
    "slide3_budget_roas",
    "slide3_strategy",
    "google_top_performer",
    "google_main_drop",
    "google_next_steps",
    "meta_top_performer",
    "meta_main_drop",
    "meta_next_steps",
    "performance_manager_narrative",
    "action_items",
}

MONTH_NUMBERS = {
    "january": "01",
    "february": "02",
    "march": "03",
    "april": "04",
    "may": "05",
    "june": "06",
    "july": "07",
    "august": "08",
    "september": "09",
    "october": "10",
    "november": "11",
    "december": "12",
}


def month_folder(year: int, month: str) -> str:
    key = month.strip().lower()
    if key not in MONTH_NUMBERS:
        raise ValueError(f"Unknown month: {month!r}. Use an English month name, e.g. March.")
    return f"{year}-{MONTH_NUMBERS[key]}"


def extract_placeholders(pptx_path: Path) -> set[str]:
    prs = Presentation(pptx_path)
    text_parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text_parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        text_parts.append(cell.text)
    return set(re.findall(r"\{\{[^}]+\}\}", "\n".join(text_parts)))


def assert_insights_shape(insights: dict) -> None:
    missing = sorted(REQUIRED_INSIGHT_KEYS - set(insights))
    if missing:
        raise ValueError(f"Claude response is missing required insight keys: {missing}")
    action_items = insights.get("action_items")
    if not isinstance(action_items, list) or len(action_items) < 5:
        raise ValueError("Claude response must include at least 5 action_items.")


def build_fake_insights(kpis: dict, client: str, month: str, year: int, next_month: str) -> dict:
    total = kpis["totals"]
    google = kpis["google"]
    meta = kpis["meta"]
    return {
        "slide3_general_insights": (
            f"{client} generated ${total['revenue']:,.2f} from ${total['cost']:,.2f} "
            f"in ad spend during {month} {year}, producing {total['roas']} ROAS."
        ),
        "slide3_budget_roas": (
            f"Google delivered {google['roas']} ROAS and Meta delivered {meta['roas']} ROAS; "
            "prioritize spend toward the stronger marginal return after reviewing funnel capacity."
        ),
        "slide3_strategy": (
            "Use this dry-run narrative only to validate the PowerPoint replacement pipeline."
        ),
        "google_top_performer": f"Google revenue was ${google['revenue']:,.2f} on ${google['cost']:,.2f} spend.",
        "google_main_drop": "Dry-run insight: review campaigns with low ROAS before scaling.",
        "google_next_steps": "Dry-run action: isolate the best funnel and validate search term quality.",
        "meta_top_performer": f"Meta revenue was ${meta['revenue']:,.2f} on ${meta['cost']:,.2f} spend.",
        "meta_main_drop": "Dry-run insight: review creative fatigue and audience overlap.",
        "meta_next_steps": "Dry-run action: refresh creative and rebalance budget by funnel.",
        "performance_manager_narrative": (
            f"Dry-run summary for {client}: validate final language with Claude before client delivery."
        ),
        "action_items": [
            f"Confirm {next_month} budget allocation.",
            "Review top ads by revenue.",
            "Check underperforming funnel stages.",
            "Validate company revenue override.",
            "Replace this dry-run copy with Claude output before final delivery.",
        ],
    }


def find_single_file(paths: Iterable[Path], label: str) -> Path:
    matches = sorted(paths)
    if not matches:
        raise FileNotFoundError(f"Missing {label}.")
    if len(matches) > 1:
        names = ", ".join(str(path) for path in matches)
        raise ValueError(f"Expected one {label}, found {len(matches)}: {names}")
    return matches[0]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Run the monthly report automation.")
    parser.add_argument("--client", required=True)
    parser.add_argument("--month", required=True)
    parser.add_argument("--year", required=True, type=int)
    parser.add_argument("--prev-month", required=True)
    parser.add_argument("--next-month", required=True)
    parser.add_argument("--company-revenue", required=True, type=float)
    parser.add_argument("--ad-revenue", required=True, type=float)
    parser.add_argument("--ad-cost", required=True, type=float)
    parser.add_argument("--media-buyer-notes", default="")
    parser.add_argument("--special-requests", default="")
    parser.add_argument(
        "--prev-data-path",
        default=None,
        help="Path to previous month's xlsx for _PREV placeholder values (e.g. data/Client/2026-02/Monthly Report February 2026.xlsx).",
    )
    parser.add_argument(
        "--use-claude",
        action="store_true",
        help="Call Claude for real insights. Omit for a no-cost local dry run.",
    )
    parser.add_argument(
        "--upload",
        action="store_true",
        help="Upload the generated PPTX to Google Drive.",
    )
    parser.add_argument(
        "--allow-untagged-template",
        action="store_true",
        help="Continue even when the template has no {{PLACEHOLDER}} tokens.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    data_dir = Path("data") / args.client / month_folder(args.year, args.month)
    xlsx_path = find_single_file(data_dir.glob("Monthly Report*.xlsx"), f"Excel export in {data_dir}")
    template_path = data_dir / "previous_report.pptx"
    if not template_path.exists():
        raise FileNotFoundError(f"Missing template: {template_path}")

    campaigns = load_sheet(xlsx_path, "Campaigns", args.client)
    ads = load_sheet(xlsx_path, "Ads", args.client)
    checkpoints = validate_or_raise(campaigns)
    kpis = build_full_kpi_report(campaigns, ads)

    prev_kpis = None
    if args.prev_data_path:
        prev_xlsx = Path(args.prev_data_path)
        if not prev_xlsx.exists():
            raise FileNotFoundError(f"Previous month xlsx not found: {prev_xlsx}")
        prev_campaigns = load_sheet(prev_xlsx, "Campaigns", args.client)
        prev_ads = load_sheet(prev_xlsx, "Ads", args.client)
        prev_kpis = build_full_kpi_report(prev_campaigns, prev_ads)

    overrides = {
        "company_revenue": args.company_revenue,
        "ad_revenue": args.ad_revenue,
        "ad_cost": args.ad_cost,
    }

    if args.use_claude:
        insights = generate_insights(
            args.client,
            args.month,
            args.year,
            args.prev_month,
            kpis,
            overrides,
            media_buyer_notes=args.media_buyer_notes,
            special_requests=args.special_requests,
        )
        assert_insights_shape(insights)
        insights_mode = "claude"
    else:
        insights = build_fake_insights(kpis, args.client, args.month, args.year, args.next_month)
        insights_mode = "dry_run_fake"

    replacements = build_replacements(
        args.client,
        args.month,
        args.year,
        args.prev_month,
        args.next_month,
        kpis,
        insights,
        overrides,
        prev_kpis=prev_kpis,
    )

    template_placeholders = extract_placeholders(template_path)
    if not template_placeholders and not args.allow_untagged_template:
        expected = "\n".join(sorted(replacements))
        raise ValueError(
            "Template has no {{PLACEHOLDER}} tokens, so the PPTX cannot be validated end-to-end yet.\n"
            f"Add the needed placeholders to {template_path}. Expected placeholders include:\n{expected}"
        )

    output_path = Path("output") / f"{args.client}_{args.month}_{args.year}_Report.pptx"
    fill_pptx(template_path, output_path, replacements)
    remaining_placeholders = extract_placeholders(output_path)

    upload_result = None
    if args.upload:
        from tools.upload_to_drive import upload_to_drive

        upload_result = upload_to_drive(output_path)

    summary = {
        "client": args.client,
        "period": f"{args.month} {args.year}",
        "xlsx_path": str(xlsx_path),
        "template_path": str(template_path),
        "output_path": str(output_path),
        "insights_mode": insights_mode,
        "checkpoints": checkpoints,
        "totals": kpis["totals"],
        "template_placeholder_count": len(template_placeholders),
        "remaining_placeholders": sorted(remaining_placeholders),
        "upload": upload_result,
    }
    print(json.dumps(summary, indent=2, default=str))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
